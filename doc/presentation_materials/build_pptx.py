# -*- coding: utf-8 -*-
"""APIEMS2026 原稿からゼミ発表用 pptx を生成する（深い青テーマ / 洗練シンプル版）。

デザイン方針:
  - テーマカラーは「深い青」の単色系。装飾を足さず余白で見せる。
  - 表紙・章扉は深い青の背景に白文字。本文は白背景。
  - 本文スライドは「タイトル → 単色の細い区切り線 → 1行キーメッセージ → コンテンツ」。
    キーメッセージは各スライドの“主張そのもの”を1文で置く。
  - 図は seminar/ の原図をそのまま使用（SVGは render_assets.py で忠実にPNG化）。

前提:  python render_assets.py を先に実行（概念図PNGを生成）。
再生成:  python build_pptx.py
"""
import os
import re as _re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# ---------------------------------------------------------------- パス
HERE = os.path.dirname(os.path.abspath(__file__))
DOC = os.path.dirname(HERE)
SEM = os.path.join(DOC, "seminar")
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)
OUT = os.path.join(HERE, "APIEMS2026_ゼミ発表.pptx")

def sem(name):  return os.path.join(SEM, name)
def asset(name): return os.path.join(ASSETS, name)

# ---------------------------------------------------------------- 配色（深い青の単色系）
INK    = RGBColor(0x26, 0x30, 0x3A)   # 本文
BLUE   = RGBColor(0x16, 0x38, 0x5C)   # 主役：深い青（見出し・線・強調・表ヘッダ）
COVER  = RGBColor(0x12, 0x31, 0x50)   # 表紙・章扉の背景
BLUESOFT = RGBColor(0x3E, 0x61, 0x8C) # 章番号などの控えめな青
LTBLUE = RGBColor(0xAF, 0xC2, 0xD8)   # 濃紺背景上のサブ文字
EN_BLUE = RGBColor(0x8C, 0xA3, 0xC0)  # 濃紺背景上の英字
MUTE   = RGBColor(0x6E, 0x7A, 0x88)   # 脚注・キャプション
FILL   = RGBColor(0xED, 0xF2, 0xF7)   # 淡い青（コールアウト・表の交互行）
ROWALT = RGBColor(0xF6, 0xF9, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HAIR   = RGBColor(0xD6, 0xDE, 0xE7)   # ごく薄い罫線

JP = "Meiryo"

# ---------------------------------------------------------------- レイアウト定数(inch)
SW, SH = 13.333, 7.5
ML = 0.62
CW = SW - 2 * ML
LINE_Y = 1.06
KEY_TOP = 1.14
CONTENT_TOP_KEY = 1.80
CONTENT_TOP_NOKEY = 1.30
CONTENT_BOTTOM = 7.02

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
_page = {"n": 0}

# ================================================================ テキスト（数式・太字）

def _set_font(run, size, color=INK, bold=False, italic=False, font=JP):
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)


def _run(paragraph, text, size, color=INK, bold=False, italic=False,
         baseline=None, font=JP):
    r = paragraph.add_run(); r.text = text
    _set_font(r, size, color=color, bold=bold, italic=italic, font=font)
    if baseline is not None:
        r._r.get_or_add_rPr().set('baseline', str(int(baseline)))
    return r


_SYM = [
    ("\\left|", "|"), ("\\right|", "|"), ("\\left(", "("), ("\\right)", ")"),
    ("\\lambda", "λ"), ("\\rho", "ρ"), ("\\delta", "δ"), ("\\Delta", "Δ"),
    ("\\beta", "β"), ("\\tau", "τ"), ("\\in", "∈"), ("\\sum", "Σ"),
    ("\\times", "×"), ("\\approx", "≈"), ("\\cdot", "·"), ("\\min", "min"),
    ("\\,", " "), ("\\;", " "), ("\\!", ""), ("\\ ", " "), ("\\{", "{"), ("\\}", "}"),
]


def _latex_clean(s):
    s = _re.sub(r'\\widehat\{([^{}]*)\}', r'\1', s)
    s = _re.sub(r'\\hat\{([^{}]*)\}', r'\1', s)
    s = _re.sub(r'\\hat\s+(\w)', r'\1', s)
    s = _re.sub(r'\\mathcal\{([^{}]*)\}', r'\1', s)
    s = _re.sub(r'\\mathrm\{([^{}]*)\}', r'\1', s)
    s = _re.sub(r'\\text\{([^{}]*)\}', r'\1', s)
    for a, b in _SYM:
        s = s.replace(a, b)
    return s


def _emit_mathexpr(paragraph, expr, size, color, bold):
    s = _latex_clean(expr)
    i, n, buf = 0, len(s), ""
    def flush():
        nonlocal buf
        if buf:
            _run(paragraph, buf, size, color=color, bold=bold, italic=True); buf = ""
    while i < n:
        c = s[i]
        if c in "_^":
            flush(); i += 1
            if i < n and s[i] == "{":
                j = s.find("}", i)
                grp = s[i + 1:j] if j != -1 else s[i + 1:]
                i = (j + 1) if j != -1 else n
            else:
                grp = s[i] if i < n else ""; i += 1
            grp = grp.strip()
            if grp:
                base = -22000 if c == "_" else 30000
                _run(paragraph, grp, size * 0.72, color=color, bold=bold,
                     italic=True, baseline=base)
        else:
            buf += c; i += 1
    flush()


def _emit_math(paragraph, seg, size, color, bold, emph):
    use = emph if (bold and emph is not None) else color
    for i, part in enumerate(seg.split("$")):
        if i % 2 == 0:
            if part:
                _run(paragraph, part, size, color=use, bold=bold)
        elif part:
            _emit_mathexpr(paragraph, part, size, use, bold)


def _add_runs(paragraph, text, size, color=INK, bold=False, emph=None):
    """**太字** と $数式$ を run に分解。emph 指定時は太字だけその色にする。"""
    state = bold
    for seg in text.split("**"):
        if seg:
            _emit_math(paragraph, seg, size, color, state, emph)
        state = not state


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    return tb, tf


def _hline(slide, x, y, w, color=BLUE, weight=1.25):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def _rect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


# ---- 作図用（角丸ボックス・コネクタ・カード） ----
CARDLINE = RGBColor(0xC9, 0xD6, 0xE5)

def _rrect(slide, x, y, w, h, fill, line=None, radius=0.10, lw=1.1):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def _boxtext(slide, x, y, w, h, lines, fill=BLUE, line=None, radius=0.10):
    sp = _rrect(slide, x, y, w, h, fill, line, radius)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (t, s, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
        _add_runs(p, t, s, color=c, bold=b, emph=c)
    return sp


def _connector(slide, x1, y1, x2, y2, color=BLUESOFT, wt=1.4, arrow=True):
    cn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(wt); cn.shadow.inherit = False
    if arrow:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def _card(slide, x, y, w, h, title, body, badge=None, accent=BLUE):
    _rrect(slide, x, y, w, h, WHITE, line=CARDLINE, radius=0.05, lw=1.1)
    tb, tf = add_textbox(slide, x + 0.2, y + 0.2, w - 0.4, 0.85)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.02
    if badge:
        _run(p, badge + " ", 18, color=accent, bold=True)
    _add_runs(p, title, 14, color=accent, bold=True, emph=accent)
    tb, tf = add_textbox(slide, x + 0.2, y + (1.12 if badge else 0.78), w - 0.4,
                         h - (1.25 if badge else 0.9))
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.1
    _add_runs(p, body, 12.5, color=INK, emph=accent)

# ================================================================ スライド枠（本文）

def _title_zone(slide, title, key):
    tb, tf = add_textbox(slide, ML, 0.34, CW, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    _add_runs(tf.paragraphs[0], title, 24.5, color=BLUE, bold=True)
    _hline(slide, ML, LINE_Y, CW, color=BLUE, weight=1.25)
    if key:
        _rect(slide, ML, KEY_TOP + 0.02, 0.065, 0.34, BLUE)   # 短い縦アクセント
        tb, tf = add_textbox(slide, ML + 0.22, KEY_TOP, CW - 0.22, 0.42,
                             anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.line_spacing = 1.0
        _add_runs(p, key, 15, color=BLUE, bold=True)


def _footer(slide, section):
    _page["n"] += 1
    if section:
        tb, tf = add_textbox(slide, ML, 7.08, 9.0, 0.30)
        _run(tf.paragraphs[0], section, 9, color=MUTE)
    tb, tf = add_textbox(slide, SW - ML - 1.2, 7.08, 1.2, 0.30)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, str(_page["n"]), 10, color=MUTE)

# ================================================================ ブロック描画（フロー）
GAP = 0.15

def _img_size(path):
    with Image.open(path) as im:
        return im.width, im.height


def render_bullets(slide, items, x, y, w):
    tb, tf = add_textbox(slide, x, y, w, 0.4)
    first, lines = True, 0
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.06
        if level == 0:
            p.space_before = Pt(7)
            _run(p, "▪ ", 12, color=BLUE, bold=True)
            _add_runs(p, text, 15.5, color=INK, emph=BLUE)
            indent = 0.0
        else:
            p.space_before = Pt(3)
            p.level = 1
            _run(p, "－ ", 12, color=MUTE)
            _add_runs(p, text, 14, color=INK, emph=BLUE)
            indent = 0.42
        cap = max(8, (w - indent) / 0.163)
        wlen = sum(2 if ord(c) > 0x2000 else 1 for c in text)
        lines += max(1, -(-wlen // int(cap)))
    return lines * 0.29 + len(items) * 0.11 + 0.05


def render_image(slide, path, x, y, w, target_h=None, target_w=None, caption=None,
                 align="center", border=False):
    iw, ih = _img_size(path)
    ar = iw / ih
    if target_h is not None:
        h = target_h; wd = h * ar
        if wd > w:
            wd = w; h = wd / ar
    else:
        wd = (target_w if target_w else w); h = wd / ar
    ix = x + (w - wd) / 2 if align == "center" else (x if align == "left" else x + (w - wd))
    pic = slide.shapes.add_picture(path, Inches(ix), Inches(y), Inches(wd), Inches(h))
    if border:
        pic.line.color.rgb = HAIR; pic.line.width = Pt(0.75)
    total = h
    if caption:
        tb, tf = add_textbox(slide, x, y + h + 0.03, w, 0.30)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _add_runs(p, caption, 10.5, color=MUTE)
        total += 0.32
    return total


def _cell_border(cell, edges=("bottom",), color="D6DEE7", w_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    tag = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}
    for e in edges:
        el = tcPr.find(qn(tag[e]))
        if el is not None:
            tcPr.remove(el)
        ln = tcPr.makeelement(qn(tag[e]), {'w': str(int(w_pt * 12700)), 'cap': 'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        clr = fill.makeelement(qn('a:srgbClr'), {'val': color})
        fill.append(clr); ln.append(fill); tcPr.append(ln)


def render_table(slide, rows, x, y, w, col_w=None, size=12.5, header=True, align=None):
    nr, nc = len(rows), len(rows[0])
    gshape = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(0.4 * nr))
    tbl = gshape.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    for sid in list(tblPr.findall(qn('a:tableStyleId'))):
        tblPr.remove(sid)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {})
    sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
    tblPr.append(sid)
    if col_w:
        tot = sum(col_w)
        for i, c in enumerate(tbl.columns):
            c.width = Inches(w * col_w[i] / tot)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(0.33)
        is_head = header and ri == 0
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if is_head else (WHITE if ri % 2 == 1 else ROWALT)
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            if align and align[ci] == 'c':
                p.alignment = PP_ALIGN.CENTER
            elif ci == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            _add_runs(p, str(val), size, color=(WHITE if is_head else INK),
                      bold=is_head, emph=(WHITE if is_head else BLUE))
            if not is_head:
                _cell_border(cell, edges=("bottom",))
    return 0.33 * nr + 0.05


def render_note(slide, text, x, y, w):
    cap = max(8, (w - 0.3) / 0.152)
    wlen = sum(2 if ord(c) > 0x2000 else 1 for c in text)
    lines = max(1, -(-wlen // int(cap)))
    h = lines * 0.27 + 0.22
    _rect(slide, x, y, w, h, FILL)
    _rect(slide, x, y, 0.06, h, BLUE)          # 左アクセントバー
    tb, tf = add_textbox(slide, x + 0.22, y, w - 0.34, h, anchor=MSO_ANCHOR.MIDDLE)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.04
    _add_runs(p, text, 13, color=INK, emph=BLUE)
    return h


def eq_png(name, latex, fontsize=22):
    path = asset(f"eq_{name}.png")
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, f"${latex}$", fontsize=fontsize, color="#16385C")
    fig.savefig(path, dpi=200, bbox_inches="tight", pad_inches=0.08, facecolor="white")
    plt.close(fig)
    return path


def render_equation(slide, latex, x, y, w, number=None, name="eq", h=0.6):
    path = eq_png(name, latex)
    iw, ih = _img_size(path); ar = iw / ih
    wd = min(w * 0.8, h * ar); hh = wd / ar
    slide.shapes.add_picture(path, Inches(x + (w - wd) / 2), Inches(y), Inches(wd), Inches(hh))
    if number:
        tb, tf = add_textbox(slide, x + w - 0.7, y + hh / 2 - 0.15, 0.7, 0.3)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        _run(p, number, 13, color=INK)
    return hh + 0.12

# ================================================================ フロー

def content_slide(title, key, section, blocks):
    slide = prs.slides.add_slide(BLANK)
    _title_zone(slide, title, key)
    _footer(slide, section)
    y0 = CONTENT_TOP_KEY if key else CONTENT_TOP_NOKEY
    _render_blocks(slide, blocks, ML, y0, CW)
    return slide


def _render_blocks(slide, blocks, x, y, w):
    cur = y
    for b in blocks:
        if 'y' in b:
            cur = b['y']
        k = b['k']
        if k == 'spacer':
            cur += b.get('h', 0.2); continue
        if k == 'bullets':
            hh = render_bullets(slide, b['items'], x, cur, b.get('w', w))
        elif k == 'image':
            hh = render_image(slide, b['path'], x, cur, b.get('w', w),
                              target_h=b.get('h'), target_w=b.get('tw'),
                              caption=b.get('caption'), align=b.get('align', 'center'),
                              border=b.get('border', False))
        elif k == 'table':
            hh = render_table(slide, b['rows'], x, cur, b.get('w', w),
                              col_w=b.get('col_w'), size=b.get('size', 12.5),
                              align=b.get('align'), header=b.get('header', True))
        elif k == 'note':
            hh = render_note(slide, b['text'], x, cur, b.get('w', w))
        elif k == 'draw':
            hh = b['fn'](slide, x, cur, b.get('w', w))
        elif k == 'eq':
            hh = render_equation(slide, b['latex'], x, cur, b.get('w', w),
                                 number=b.get('number'), name=b.get('name', 'eq'),
                                 h=b.get('h', 0.6))
        elif k == 'row':
            ratio = b.get('ratio', [1] * len(b['cols']))
            tot = sum(ratio); gap = b.get('gap', 0.34)
            avail = w - gap * (len(b['cols']) - 1)
            hmax, cx = 0, x
            for col, rr in zip(b['cols'], ratio):
                cwid = avail * rr / tot
                hmax = max(hmax, _render_blocks(slide, col, cx, cur, cwid))
                cx += cwid + gap
            hh = hmax
        else:
            hh = 0
        cur += hh + GAP
    return cur - y

# ================================================================ 表紙・章扉

def title_slide():
    slide = prs.slides.add_slide(BLANK)
    _rect(slide, 0, 0, SW, SH, COVER)
    # 上部の細いアクセント
    _hline(slide, 0.95, 2.28, 2.6, color=BLUESOFT, weight=2.0)
    tb, tf = add_textbox(slide, 0.95, 2.45, SW - 1.9, 1.9)
    p = tf.paragraphs[0]; p.line_spacing = 1.12
    _run(p, "ジョブショップ再スケジューリングにおける", 31, color=WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.12
    _run(p2, "安定性誘導演算子の非対称効果", 31, color=WHITE, bold=True)
    p3 = tf.add_paragraph(); p3.space_before = Pt(10)
    _run(p3, "— 軌道探索と集団探索の横断比較 —", 17, color=LTBLUE, bold=True)
    # 英題
    tb, tf = add_textbox(slide, 0.97, 4.55, SW - 1.9, 0.7)
    p = tf.paragraphs[0]; p.line_spacing = 1.05
    _run(p, "Asymmetric Effects of Stability-Inducing Operators across "
            "Trajectory and Population Search in Job-Shop Rescheduling",
         12.5, color=EN_BLUE, italic=True)
    # 著者
    _hline(slide, 0.97, 5.55, 3.4, color=BLUESOFT, weight=1.0)
    tb, tf = add_textbox(slide, 0.97, 5.68, SW - 1.9, 1.3)
    rows = [("早稲田大学 創造理工学部 経営システム工学科", 14.5, LTBLUE, False),
            ("鬼頭 拓海", 18, WHITE, True),
            ("ゼミ発表　2026-07-21　｜　投稿先：APIEMS 2026（釜山）", 12, EN_BLUE, False)]
    for i, (txt, sz, col, bold) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5 if i else 0)
        _run(p, txt, sz, color=col, bold=bold)
    return slide


def divider_slide(no, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    _rect(slide, 0, 0, SW, SH, COVER)
    tb, tf = add_textbox(slide, 1.15, 2.35, 1.7, 1.5, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], no, 78, color=BLUESOFT, bold=True)
    tb, tf = add_textbox(slide, 2.75, 2.35, SW - 3.6, 1.5, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _run(p, title, 30, color=WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph(); p2.space_before = Pt(10)
        _run(p2, subtitle, 14, color=LTBLUE)
    _hline(slide, 2.78, 4.28, 5.2, color=BLUESOFT, weight=1.25)
    return slide

# ================================================================ 作図（1〜2章のビジュアル）

def diag_taxonomy(slide, x, y, w):
    """1.1：外乱対応の分類フロー（静的/動的 → 予測リアクティブ → 作業遅延）"""
    bh = 0.60; sy = y + 0.10; dy = y + 0.98
    _boxtext(slide, x + 0.05, y + 0.34, 1.85, 0.90, [("外乱への\n対応", 13, WHITE, True)],
             fill=BLUE, radius=0.12)
    _boxtext(slide, x + 2.45, sy, 3.0, bh, [("静的（事前耐性設計）", 12, MUTE, False)],
             fill=FILL, line=CARDLINE)
    _boxtext(slide, x + 2.45, dy, 3.0, bh, [("動的（事後修正）＝再スケジューリング", 11, BLUE, True)],
             fill=FILL, line=BLUE)
    _boxtext(slide, x + 5.85, dy, 3.1, bh, [("予測リアクティブ再スケジューリング", 11.5, WHITE, True)],
             fill=BLUE)
    _boxtext(slide, x + 9.35, dy, 2.6, bh, [("作業遅延に着目", 12.5, BLUE, True)],
             fill=WHITE, line=BLUE)
    _boxtext(slide, x + 6.55, dy - 0.34, 1.7, 0.32, [("★ 本研究の対象", 10.5, WHITE, True)],
             fill=BLUESOFT, radius=0.30)
    _connector(slide, x + 1.90, y + 0.62, x + 2.45, sy + bh / 2, arrow=False)
    _connector(slide, x + 1.90, y + 1.08, x + 2.45, dy + bh / 2, arrow=False)
    _connector(slide, x + 5.45, dy + bh / 2, x + 5.85, dy + bh / 2)
    _connector(slide, x + 8.95, dy + bh / 2, x + 9.35, dy + bh / 2)
    return 1.78


def diag_tradeoff(slide, x, y, w):
    """1.2：効率⇄安定のトレードオフ ＋ 変更コストのチップ"""
    bh = 0.88
    _boxtext(slide, x + 0.15, y + 0.05, 4.55, bh,
             [("効率（メイクスパン MS）", 13, BLUE, True), ("を最小化", 10.5, INK, False)],
             fill=FILL, line=BLUE)
    _boxtext(slide, x + w - 4.70, y + 0.05, 4.55, bh,
             [("安定性（$S_p$ からの変更量 $D$）", 12.5, BLUE, True), ("を最小化", 10.5, INK, False)],
             fill=FILL, line=BLUE)
    cx = x + w / 2
    tb, tf = add_textbox(slide, cx - 0.85, y + 0.02, 1.7, bh)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "⇄", 24, color=BLUE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, "トレードオフ", 11, color=BLUE, bold=True)
    ty = y + 1.22
    tb, tf = add_textbox(slide, x + 0.05, ty, 5.5, 0.32)
    _add_runs(tf.paragraphs[0], "大幅な変更が生むコスト", 12.5, color=INK, bold=True, emph=BLUE)
    chips = ["現場の混乱", "段取り替え", "資材・治具の再手配", "作業者の再配置", "下流・外注への波及"]
    n = len(chips); cwid = (w - 0.10 - 0.14 * (n - 1)) / n; cyy = ty + 0.40
    for i, t in enumerate(chips):
        _boxtext(slide, x + 0.05 + i * (cwid + 0.14), cyy, cwid, 0.48,
                 [(t, 11, BLUE, False)], fill=FILL, line=CARDLINE, radius=0.20)
    return 2.28


def diag_obj_cards(slide, x, y, w):
    """1.3：研究目的の3カード"""
    cards = [
        ("①", "特性と探索構造の適合性分析",
         "高品質初期解 $S_p$ の存在が、軌道／集団の探索挙動をどう変えるかを統制下で分析する。"),
        ("②", "安定性誘導演算子の提案",
         "解を $S_p$ へ能動的に引き寄せる演算子（PR・repair）を設計し、ベース構造との相互作用を分析する。"),
        ("③", "多角的評価方法論の構築",
         "Pareto 覆域・安定性帯別・収束速度を統合した評価フレームワークを構築する。"),
    ]
    gap = 0.4; cw = (w - 2 * gap) / 3; ch = 2.45
    for i, (bd, ti, bo) in enumerate(cards):
        _card(slide, x + i * (cw + gap), y + 0.05, cw, ch, ti, bo, badge=bd)
    return ch + 0.1


def diag_hyp(slide, x, y, w):
    """H1/H2 の2カード"""
    gap = 0.4; cw = (w - gap) / 2; ch = 1.35
    _card(slide, x, y + 0.05, cw, ch, "H1（適合性）",
          "軌道 ILS は $S_p$ 近傍を充填的に覆う。集団 Memetic は交叉で $S_p$ から飛び、近傍が粗い。")
    _card(slide, x + cw + gap, y + 0.05, cw, ch, "H2（補完）",
          "$S_p$ 誘導演算子(PR・repair)は集団の弱点を補うが、自力充填済みの ILS は伸びしろが小さい ＝ ホスト依存で非対称。")
    return ch + 0.1


# ================================================================ 本編
def build():
    # 1. 表紙
    title_slide()

    # 2. 研究概要
    content_slide(
        "研究概要", "同一の安定性誘導演算子が、載せる探索構造しだいで効き方を変える——その非対称性を統制下で切り分ける",
        "Overview",
        [
            {'k': 'bullets', 'items': [
                (0, "**予測リアクティブ再スケジューリング**では、外乱前の高品質スケジュール **$S_p$** が既に存在し、修正解には効率（メイクスパン）と安定性（$S_p$ からの変更量の小ささ）が同時に求められる。"),
                (0, "→ 探索の勝負は「大域探索力」ではなく、**$S_p$ 近傍（高安定領域）をどれだけ良く充填できるか**という軸を新たに帯びる。"),
                (1, "**H1（適合性）**：軌道ベース(ILS)は連続変形でこの近傍を自力充填。集団ベース(Memetic)は交叉ゆえ充填が構造的に粗い。"),
                (1, "**H2（補完）**：解を $S_p$ へ引き寄せる**安定性誘導演算子(PR・repair)**は集団の弱点を補うため、効果が**ホスト依存で非対称**に現れる。"),
                (0, "主眼は新手法の性能競争ではなく、**演算子とホスト構造の相互作用を統制下で解明**すること（8 シナリオ × 7 手法 × n=10）。"),
            ]},
        ])

    # 3. 目次
    content_slide(
        "目次", None, "Contents",
        [
            {'k': 'bullets', 'items': [
                (0, "**1. 研究背景・目的** ｜ 外乱・再スケジューリング・安定性のトレードオフ"),
                (0, "**2. 既存研究と本研究の位置づけ** ｜ 研究ギャップと 2 仮説（H1 / H2）"),
                (0, "**3. 問題設定と提案手法** ｜ 問題定義・安定性誘導演算子（PR / repair）・評価フレームワーク"),
                (0, "**4. 計算機実験** ｜ H1（適合性）・H2（非対称効果）・総合スコアボード"),
                (0, "**5. 結論** ｜ 相補構造・限界・今後の課題"),
            ]},
        ])

    # ===== 1
    divider_slide("1", "研究背景・目的", "外乱・再スケジューリング・安定性のトレードオフ")

    content_slide(
        "1.1 外乱と再スケジューリング", "本研究の対象を「予測リアクティブ再スケジューリング × 作業遅延」に位置づける",
        "1. 研究背景・目的",
        [
            {'k': 'bullets', 'items': [
                (0, "生産現場では**作業遅延・機械故障**などの外乱が頻発し、当初スケジュールの実行を困難にする。対応は静的・動的の 2 系統に大別される。"),
            ]},
            {'k': 'draw', 'fn': diag_taxonomy},
            {'k': 'bullets', 'items': [
                (0, "**作業遅延**に着目する理由：機械割当を保ち、処理**順序の再調整**だけで対処できる ＝ **作業集合を変えない**。"),
                (1, "ゆえに**修正解を $S_p$ 近傍に保つ**という本研究の前提が、最も素直に成立する外乱クラスである。"),
            ]},
        ])

    content_slide(
        "1.2 効率と安定性のトレードオフ", "変更コストゆえ安定性は効率(MS)と並ぶ目的——効率と安定性の多目的最適化として定式化する",
        "1. 研究背景・目的",
        [
            {'k': 'draw', 'fn': diag_tradeoff},
            {'k': 'bullets', 'items': [
                (0, "→ 安定性は効率(MS)と並ぶ目的であり、**効率性と安定性の多目的最適化**として定式化する。"),
            ]},
            {'k': 'note', 'text': "再スケジューリングの本質的特殊性：**高品質初期解 $S_p$ が既に存在**し、最適解は $S_p$ 近傍に分布しやすい。ゆえに手法の優劣は大域探索力だけでは決まらない。"},
        ])

    content_slide(
        "1.3 研究目的", "$S_p$ が既にある特性を軸に、探索構造との適合性・安定性誘導演算子・多角的評価の 3 点を目的に据える",
        "1. 研究背景・目的",
        [
            {'k': 'draw', 'fn': diag_obj_cards},
            {'k': 'note', 'text': "主眼は最先端手法との性能競争ではなく、**演算子とホスト構造の相互作用の構造分析**にある。"},
        ])

    # ===== 2
    divider_slide("2", "既存研究と本研究の位置づけ", "研究ギャップと 2 仮説（H1 / H2）")

    content_slide(
        "2.1 研究ギャップと本研究の位置づけ", "既存に残る 3 つの課題に、本研究は「統制比較・演算子・多軸評価」で応える",
        "2. 既存研究と位置づけ",
        [
            {'k': 'bullets', 'items': [
                (0, "効率と安定性の同時最適化は Wu ら[14]以降の中心課題。安定性を**範囲の制限**で担う系譜(match-up[15]・AOR[16]・スコープ[27])もあるが、次の **3 つの課題**が残る。"),
            ]},
            {'k': 'table', 'rows': [
                ["残された課題", "既存研究", "本研究のアプローチ"],
                ["① 適合性分析の欠如", "GA 中心。$S_p$ が既にある特性と軌道／集団の適合性を正面分析した例がない",
                 "N5 を揃えた軌道(ILS)と集団(Memetic)の**統制比較**【H1】"],
                ["② 演算子としての機構の欠如", "安定性の確保は全て**範囲限定**（探索空間の制限）",
                 "解を $S_p$ へ引き寄せる**安定性誘導演算子**(PR・repair)【H2】"],
                ["③ 評価方法論の不足", "単一重みのスカラー比較が主",
                 "**3 指標**（Pareto 覆域・安定帯別・速度）で評価"],
            ], 'col_w': [0.82, 1.46, 1.72], 'size': 11.5},
            {'k': 'draw', 'fn': diag_hyp},
        ])

    # ===== 3
    divider_slide("3", "問題設定と提案手法", "問題定義・安定性誘導演算子（PR / repair）・評価フレームワーク")

    content_slide(
        "3.1 問題定義 — 再スケジューリングの力学", "凍結部を固定し、$t_r$ 以降の「機械ごとの処理順序」だけを動かす順列最適化に帰着する",
        "3. 問題設定と提案手法",
        [
            {'k': 'bullets', 'items': [
                (0, "**元スケジュール $S_p$**：静的 JSSP の高品質解（全手法の共通入力）。**外乱**は単一の作業遅延（遅延量 $\\Delta$）。"),
                (0, "**right-shift 解 $S_{RSR}$**：機械上の順序を保ち遅延を吸収した現場の実スケジュール（$D=0$ に対応）。"),
                (0, "**凍結**：再スケ時刻 $t_r$ 前に開始済みの作業は凍結し、$t_r$ 以降を**最適化対象 $\\mathcal{O}_{opt}$** とする。"),
                (0, "**決定変数**：機械割当は固定 → 対象作業の**機械ごとの処理順序**のみ。"),
            ]},
            {'k': 'bullets', 'items': [
                (0, "**安定性指標（順列偏差）**：探索機構(swap・N5)が動く順列空間と一致し、MS と独立性が高い。"),
            ]},
            {'k': 'eq', 'name': 'stab',
             'latex': r'D(S_p,S_q)=\sum_{(i,j)\in\mathcal{O}_{opt}}\left|\,r_{ij}^{\,p}-r_{ij}^{\,q}\,\right|',
             'number': '(1)', 'h': 0.58},
            {'k': 'bullets', 'items': [
                (1, "$r_{ij}$＝機械 $i$ 上のジョブ $j$ の処理順位。値が小さいほど高安定（$D=0$＝順序不変の $S_{RSR}$）。"),
            ]},
        ])

    content_slide(
        "3.1 多目的とスカラー化", "重み付き和で掃引——運用の選好設定と整合し、軌道と集団に同一目的を共有させる「統制の共通土台」になる",
        "3. 問題設定と提案手法",
        [
            {'k': 'bullets', 'items': [
                (0, "$\\min_{S_q}(MS(S_q),\\, D(S_p,S_q))$ を、重み $\\lambda\\in[0,1]$ の**重み付き和**で解く（$\\hat\\cdot$ は min-max 正規化）。"),
            ]},
            {'k': 'eq', 'name': 'scal',
             'latex': r'F(S_q)=\lambda\,\hat D(S_p,S_q)+(1-\lambda)\,\widehat{MS}(S_q)',
             'number': '(2)', 'h': 0.56},
            {'k': 'bullets', 'items': [
                (0, "**なぜ Pareto-native 手法でなく重みスカラー化か**"),
                (1, "**運用形態との整合**：効率と安定の優先度は意思決定者が事前に重みとして与える量（再スケの標準枠組[3]）。"),
                (1, "**統制比較の共通基盤**：単一解 ILS と集団 Memetic に同一目的 $F(\\lambda)$ を共有させ、同一機構を無改変で載せ**構造差だけ**を切り出せる。"),
                (0, "単一重み依存を避け $\\lambda$ を複数点掃引し、**UEA[9]** で全非劣解を統合して評価する。"),
            ]},
        ])

    content_slide(
        "3.2 ベース探索構造と H1 の機序", "軌道は摂動強度＝$S_p$ からの距離を制御でき近傍を充填、集団は破壊的な交叉で子個体が $S_p$ から飛散する（H1）",
        "3. 問題設定と提案手法",
        [
            {'k': 'row', 'ratio': [1.05, 1], 'cols': [
                [{'k': 'image', 'path': asset("concept_2struct.png"), 'h': 2.95,
                  'caption': "集団は交叉で子個体が $S_p$ から飛散（左）／軌道は連続変形で $S_p$ 近傍を充填（右）"}],
                [{'k': 'bullets', 'items': [
                    (0, "**軌道(ILS)**：深掘り(局所探索)と脱出(摂動)が分離 → **摂動強度＝$S_p$ からの距離**を直接制御でき、$S_p$ 起点に連続変形で近傍を充填。"),
                    (0, "**集団(Memetic)**：GA に N5 を Lamarckian 適用。だが**交叉は 2 親を切り貼りする破壊的操作**で、子個体が $S_p$ 近傍から飛散する。"),
                    (0, "**局所探索 N5**[1] は両者で共有（統制）→ 高安定充填の差は**交叉の有無＝探索構造**だけに帰属。"),
                    (0, "∴ **H1**：同じ N5 でも集団は高安定領域の充填が構造的に粗い。"),
                ]}],
            ]},
            {'k': 'note', 'text': "集団が $S_p$ から遠い解を多数保つこの性質は、次節の**安定性誘導演算子（PR・repair）が働く土壌**になる（→ H2）。"},
        ])

    content_slide(
        "3.3 安定性誘導演算子 — PR・repair", "guiding solution を $S_p$ に固定＝解を「安定性アンカー」へ方向づける移動。同一演算子が両ホストに載る",
        "3. 問題設定と提案手法",
        [
            {'k': 'row', 'ratio': [1.05, 1], 'cols': [
                [{'k': 'image', 'path': asset("concept_direct_swap.png"), 'h': 3.1}],
                [{'k': 'bullets', 'items': [
                    (0, "**PR（Path Relinking）**[5]：現在解→$S_p$ の経路を swap で辿り経路上最良解を返す。評価回数を **$O(d)$**（$d$=不一致数）に抑える。"),
                    (0, "**repair（Mini-PR kick）**：PR を数手で打ち切り、$S_p$ 方向へ引き戻してから局所探索で再最適化。深さで安定側フロントを面で覆う。"),
                    (0, "**両ホストに載るのは同一演算子**。異なるのは発動法だけ：軌道＝停滞時の摂動／集団＝個体ごとの精緻化。"),
                ]}],
            ]},
        ])

    content_slide(
        "3.4 評価フレームワーク — 3 指標", "単一指標では相補構造が見えない——「総合品質・高安定充填・速度」の 3 問を 3 指標で切り分ける",
        "3. 問題設定と提案手法",
        [
            {'k': 'row', 'ratio': [1.15, 1], 'cols': [
                [{'k': 'table', 'rows': [
                    ["指標", "意味", "役割"],
                    ["統合 HV", "全領域の hypervolume", "総合品質"],
                    ["高安定 HV", "$D<$P50（$S_p$ 近傍）限定", "本命（充填度）"],
                    ["AOC", "HV-対-log 時間曲線の時間平均", "アンタイム（速さ）"],
                ], 'col_w': [0.8, 1.55, 0.95], 'size': 11.5, 'align': [None, None, 'c']}],
                [{'k': 'image', 'path': sem("core_v3_metrics_explained.png"), 'h': 2.15}],
            ]},
            {'k': 'bullets', 'items': [
                (0, "HV は各シナリオで $[0,1]^2$ 正規化 ＋ 参照点 $(1.1,1.1)$。AOC は壁時計・対数時間で手法間 apples-to-apples。"),
                (0, "**統計**：各シナリオ内で片側 Wilcoxon ＋ Cliff's $\\delta$。横断は Friedman 平均順位 ＋ Kendall's $W$（相関シナリオゆえ探索的要約）。$|\\delta|$=1.0 は完全分離（$p\\approx0.001$）。"),
            ]},
        ])

    # ===== 4
    divider_slide("4", "計算機実験", "H1（適合性）・H2（非対称効果）・総合スコアボード")

    content_slide(
        "4.1 実験設定", "同一 $S_p$ のまま再スケ率 ρ だけを段階変化させ、手法差の ρ 依存を交絡なく切り分ける設計",
        "4. 計算機実験",
        [
            {'k': 'table', 'rows': [
                ["対象 7 手法", "baseline", "＋ repair", "＋ PR"],
                ["軌道（ILS）", "ILS-baseline", "ILS+repair", "ILS+PR"],
                ["集団（Memetic）", "Memetic-LS", "Memetic+repair", "Memetic+PR"],
            ], 'col_w': [1.05, 1.0, 1.0, 1.0], 'size': 12, 'align': [None, 'c', 'c', 'c']},
            {'k': 'bullets', 'items': [
                (1, "＋ **GA**（局所探索なしの参考ベースライン）で計 **7 手法**。$S_p$（GA-500 の高品質解）を全手法で共有。"),
            ]},
            {'k': 'row', 'ratio': [1, 1.05], 'cols': [
                [{'k': 'table', 'rows': [
                    ["項目", "設定"],
                    ["インスタンス", "mt10・la21・la36・la40・ta21"],
                    ["重み λ", "10 点（0〜0.9・0.1 刻み）"],
                    ["Trial", "n=10（解析単位＝1 trial）"],
                    ["計算予算", "ILS 3000 反復／GA・Memetic 500 世代"],
                ], 'col_w': [0.8, 1.7], 'size': 11}],
                [{'k': 'bullets', 'items': [
                    (0, "各シナリオの%は**再スケ率** $\\rho=n_{res}/\\text{ops}$（統制軸）：la36 27/54/73・ta21 32/82・mt10 72・la21 35・la40 32。"),
                    (0, "**la36 ラダー**・**ta21 対**は同一 $S_p$ で $\\rho$ だけ段階変化 → $\\rho$ 依存を交絡なく分離。"),
                    (0, "外乱は単一作業の完了遅延 $\\Delta$=60〜148。環境：Ryzen 5 7530U／Python 3.12・全乱数シード固定。"),
                ]}],
            ]},
        ])

    content_slide(
        "4.2 結果1：軌道(ILS) vs 集団(Memetic)【H1】", "同じ N5 を積んでも、高安定領域は ILS が全 8 シナリオで完全分離——差は局所探索でなく探索構造",
        "4. 計算機実験",
        [
            {'k': 'image', 'path': sem("core_v3_claim1.png"), 'h': 3.45,
             'caption': "局所探索(N5)を揃えた ILS-baseline vs Memetic-LS ｜ 統合HV=互角／高安定HV=ILS完全優越／AOC=6/8でILS優位"},
            {'k': 'bullets', 'items': [
                (0, "**高安定 HV：ILS が全 8 シナリオを完全分離で上回る**（$p$=0.001, $|\\delta|$=1.0）。低 $\\rho$ の 3 シナリオでは Memetic は高安定域に**1 解も到達できず**、残る 5 でも ILS が 2〜4.5 倍。Holm 補正後も全シナリオ有意。"),
                (0, "**統合 HV は互角**（ILS 5 勝・Memetic 3 勝、勝者は $\\rho$ で入れ替わる）。差は**局所探索の有無でなく探索構造**に由来する。"),
            ]},
        ])

    content_slide(
        "4.2 H1 の構造的原因 — 訪問密度差", "ILS は低 $D$ フロント帯に集中、集団は交叉ゆえ低 $D$ でも MS を煮詰めきれず充填できない",
        "4. 計算機実験",
        [
            {'k': 'row', 'ratio': [1.22, 1], 'cols': [
                [{'k': 'image', 'path': sem("h1_density.png"), 'h': 3.95}],
                [{'k': 'bullets', 'items': [
                    (0, "訪問密度差マップ（赤＝ILS 密／青＝Memetic 密、手法ごと総和 1 で正規化）。"),
                    (0, "**ILS は低 $D$ のフロント帯**（$S_p$ 近傍で良 MS を保てる帯）に集中。"),
                    (0, "**Memetic は高 $D$ の不安定領域に分散** ＝ 同じ低 $D$ でも交叉ゆえ MS を煮詰めきれず Pareto 的に充填できない。"),
                    (0, "→ これが高安定 HV 差の**構造的原因**。集団が $S_p$ から遠い解を多数持つ性質は、次の H2 の土壌になる。"),
                ]}],
            ]},
        ])

    content_slide(
        "4.3 結果2：PR・repair の非対称効果【H2】", "同一演算子が集団の高安定 HV を 2 倍超に引き上げ ILS 水準へ。だが自力充填済みの軌道はほぼ頭打ち＝非対称",
        "4. 計算機実験",
        [
            {'k': 'image', 'path': sem("core_v3_claim2.png"), 'h': 3.45,
             'caption': "baseline に機構を追加したときの高安定 HV の伸び（ホスト別）"},
            {'k': 'bullets', 'items': [
                (0, "**集団(Memetic)：高安定 HV を大幅改善**（全 8 シナリオで 2 倍以上, $p$=0.001, $|\\delta|$=1.0）。届かなかった $S_p$ 近傍を機構が直接充填 → **高安定で ILS に追いつき、統合 HV では追い越す**。"),
                (0, "**軌道(ILS)：ほぼ頭打ち**。baseline 時点で自力充填済み。機構が有意なのは最高 $\\rho$ 帯の la36L(73%)・ta21L(82%)のみ（Holm 後は ta21L のみ残存, $p_{adj}$≈0.008）。"),
            ]},
        ])

    content_slide(
        "4.3 H2 の機構的原因 — PR 経路統計", "集団は経路が長く道中で改善を 30〜65% 発見、ILS は経路が短く発見ほぼ 0%——方向づけは軌道では空振り",
        "4. 計算機実験",
        [
            {'k': 'row', 'ratio': [1.2, 1], 'cols': [
                [{'k': 'image', 'path': sem("mech_pr.png"), 'h': 3.7}],
                [{'k': 'bullets', 'items': [
                    (0, "**Memetic**：経路長 $d_0$ が大きく、経路上で**約 30〜65%** の確率で改善解を発見。"),
                    (0, "**ILS**：$S_p$ 近傍に張り付き経路が短く、改善発見率は**全シナリオほぼ 0%**（ta21L でも 0.4%）。"),
                    (0, "方向づけ移動は ILS では空振り。la36L・ta21L で僅かに効くのは、**キック後の局所探索**が埋め残しを再最適化する寄与による。"),
                ]}],
            ]},
        ])

    content_slide(
        "4.3 PR か repair か — ホスト依存の使い分け", "使い分けが要るのは集団だけ：最終品質なら PR、アンタイム重視なら repair（軌道はどちらでも同じ）",
        "4. 計算機実験",
        [
            {'k': 'table', 'rows': [
                ["ホスト", "統合／高安定 HV（品質）", "AOC（アンタイム）", "使い分け"],
                ["ILS（$d$ 小）", "機構を足しても頭打ち＝タイ", "全 8 でタイ（発動が停滞時・経路短く $O(d)$）", "選択は問題にならない"],
                ["Memetic（$d$ 大）", "PR が僅かに優る", "repair が全 8 で PR に優る（$\\delta$ 最大 +1.0）", "品質＝PR ／ 予算厳＝repair"],
            ], 'col_w': [0.78, 1.35, 1.55, 1.05], 'size': 11.5},
            {'k': 'bullets', 'items': [
                (0, "**PR** は $S_p$ までの長い経路を辿り切ってから最良中間解を返すため**立ち上がりが遅い**。**repair** は数手で打ち切り即座に incumbent を更新するため**アンタイムに強い**。"),
                (0, "→ 最終品質最優先なら **Memetic+PR**、計算予算が厳しくアンタイム重視なら **repair**。ILS では既定 repair で足りる。"),
            ]},
        ])

    content_slide(
        "4.4 総合スコアボード — 指標で首位が替わる", "評価軸が変われば首位も替わる：総合品質＝Memetic+機構／安定・速度＝ILS系。万能手法は存在しない",
        "4. 計算機実験",
        [
            {'k': 'row', 'ratio': [1, 1, 1], 'cols': [
                [{'k': 'image', 'path': sem("scoreboard_union.png"), 'h': 2.5,
                  'caption': "(a) 統合HV — 首位 Memetic+PR"}],
                [{'k': 'image', 'path': sem("scoreboard_highstab.png"), 'h': 2.5,
                  'caption': "(b) 高安定HV — 首位 ILS 系"}],
                [{'k': 'image', 'path': sem("scoreboard_aoc.png"), 'h': 2.5,
                  'caption': "(c) AOC — 首位 ILS 系"}],
            ], 'gap': 0.25},
            {'k': 'bullets', 'items': [
                (0, "Friedman 平均順位は 3 指標とも明確に分離（Kendall's $W$=0.59／0.81／0.63, $p$<0.0001, 相関ゆえ探索的要約）。"),
                (0, "**機構なしの素の集団(GA・Memetic-LS)だけが高安定 HV で ARPD≈70〜78% と壊滅**し、$S_p$ 近傍に届かず二分される（H1・H2 の予言どおり）。"),
            ]},
        ])

    content_slide(
        "4.4 結果の統合 — 相補構造と閾値頑健性", "「素の集団だけが高安定で壊滅」する二分構造は閾値 P25〜P75 で不変——相補構造は頑健",
        "4. 計算機実験",
        [
            {'k': 'bullets', 'items': [
                (0, "**統合 HV（品質）**：機構込み集団が首位群（Memetic+PR 2.0・+repair 2.5）。素の集団と互角だった ILS の**上に**機構込み Memetic が立つ（首位は leave-one-out に頑健）。"),
                (0, "**高安定 HV（本命）**：ILS 系と機構込み Memetic が首位群(2.6〜3.4)、素の集団だけ壊滅（GA 6.4・Memetic-LS 6.6）。"),
                (0, "**AOC（速さ）**：ILS 系 3 種が首位群(2.5〜2.8)。ウォームアップの遅い Memetic+PR・GA が下位。"),
                (0, "**閾値頑健性**：高安定 HV の P50 を P25〜P75 で掃引しても「素の集団＝最下位群／ILS系・機構込Memetic＝上位群」の**二分構造は不変**。動くのは首位群内の僅差順位のみ。"),
            ]},
            {'k': 'note', 'text': "**発散型 ILS**（$S_p$ 起点に外へ）と**収束型 Memetic+PR**（散った集団を $S_p$ へ）は逆向きだが類似の最終 Pareto に到達。ILS は早期から良い incumbent を持ち、AOC でこの交差が ILS 優位として現れる。"},
        ])

    content_slide(
        "4.4【探索的】統合 HV の勝者と再スケ率", "総合HVの勝者は ρ≈50% で二分——ただし交絡ゆえ見立て。本命の結論はこれに依存しない",
        "4. 計算機実験",
        [
            {'k': 'bullets', 'items': [
                (0, "統合 HV の勝者は $\\rho$ と対応し、**~50% を境に低で ILS・中〜高で Memetic** に二分（同一 la36 ラダー 27/54/73% で ILS→Memetic→Memetic）。可動工程の絶対数や問題規模では捉えられない。"),
                (0, "**機序の見立て**：安定性項は $S_p$($D$=0) を引力点とする。$\\rho$ 小では良解が $S_p$ 近傍に限られ ILS の近傍充填で足りる。$\\rho$ 増で効率端が遠のくと、単一軌道の ILS は引力圏を脱しにくく、$S_p$ から離れた個体を保つ Memetic だけが届く（H1 と同型）。"),
            ]},
            {'k': 'note', 'text': "ただし ta21L(82%)は例外（≒タイ, $p$=0.053）で、$S_p$ 品質の交絡と順列偏差表現への依存も残るため**探索的見立て**に留める。本命の高安定 HV 優位・機構非対称はこの対応に依存せず、全 8 シナリオの直接検定で独立に確立している。"},
        ])

    # ===== 5
    divider_slide("5", "結論", "相補構造・限界・今後の課題")

    content_slide(
        "5. 結論 — 3 つの主張", "安定性を「演算子」化して同一機構を両ホストへ移植し、そのホスト依存の非対称性を統制下で切り分けた",
        "5. 結論",
        [
            {'k': 'bullets', 'items': [
                (0, "**① 軌道(ILS)は高安定領域を効率的に充填する【H1】**：総合品質は互角だが、本命の高安定領域で同一 N5 の集団を全 8 シナリオで完全分離で上回る（$p$=0.001, $|\\delta|$=1.0）。局所探索の有無でなく**探索構造**に由来。"),
                (0, "**② PR・repair の効果はホスト依存で非対称【H2】**：集団の高安定 HV を 2 倍以上押し上げ ILS 水準へ。自力充填済みの軌道はほぼ頭打ち（ただし極端な高 $\\rho$=ta21L 82% では軌道にも有意に残る）。"),
                (0, "**③ 軌道と集団の相補構造**：総合品質＝Memetic+機構／安定性＝ILS と Memetic+機構の双方／速度＝ILS 系。**評価指標により最良手法が替わる**——両構造は再スケの異なる要求を補い合う。"),
            ]},
            {'k': 'note', 'text': "**中心的貢献**：安定性レバーを演算子として実装したことで同一機構を両ホストへ移植でき、そのホスト依存の非対称性を統制下で切り分けたこと。"},
        ])

    content_slide(
        "5. 限界と今後の課題", "主結論は「作業集合を保つ作業遅延・順列偏差」で確立——集合を変える外乱や開始時刻偏差が次の射程",
        "5. 結論",
        [
            {'k': 'row', 'ratio': [1, 1], 'cols': [
                [{'k': 'bullets', 'items': [
                    (0, "**限界**"),
                    (1, "n=10 に基づく（飽和した主結論は頑健、境界事例は別）。"),
                    (1, "安定性を**順列偏差でのみ**測定。開始時刻偏差下での妥当性は未検証。"),
                    (1, "H1・H2 は**作業集合を保つ作業遅延**で確立。作業集合を変える外乱では $S_p$ が完全な参照解でなくなり、式(1)の定義自体が変わる。"),
                ]}],
                [{'k': 'bullets', 'items': [
                    (0, "**今後の課題**"),
                    (1, "開始時刻偏差下での再検証。"),
                    (1, "特急ジョブ割込み・ジョブ削除など**作業集合を変える外乱**への拡張。"),
                    (1, "範囲レバーと演算子レバーの統合（波及範囲内での誘導演算子併用）。"),
                    (1, "集団側の代替処方（$S_p$ 偏向交叉）との比較、Pareto-native 手法(NSGA-II 等)での再検証。"),
                ]}],
            ]},
        ])

    prs.save(OUT)
    print("saved:", OUT, " slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
